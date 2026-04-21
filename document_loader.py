"""
document_loader.py
PDF / PPT 파일에서 텍스트를 추출하는 모듈.
사내 규정처럼 조항 구조가 있는 문서를 전처리합니다.
"""

import re
from pathlib import Path
from typing import List
from langchain_core.documents import Document


# ─────────────────────────────────────────
# PDF 로더
# ─────────────────────────────────────────
def load_pdf(file_path: str) -> List[Document]:
    """
    pymupdf(fitz)로 PDF 텍스트 추출.
    페이지별로 Document 객체를 생성하고 메타데이터를 붙입니다.
    """
    import fitz  # pymupdf

    docs = []
    pdf = fitz.open(file_path)

    for page_num, page in enumerate(pdf, start=1):
        text = page.get_text("text").strip()
        if not text:
            continue

        docs.append(Document(
            page_content=text,
            metadata={
                "source": Path(file_path).name,
                "page": page_num,
                "file_type": "pdf",
            }
        ))

    pdf.close()
    print(f"[PDF] {Path(file_path).name} → {len(docs)} 페이지 추출")
    return docs


# ─────────────────────────────────────────
# Markdown 로더
# ─────────────────────────────────────────
def load_markdown(file_path: str) -> List[Document]:
    """
    마크다운 파일을 헤더(#, ##, ###) 기준으로 섹션 분할하여 로드합니다.
    헤더가 없으면 전체를 하나의 Document로 반환합니다.
    """
    HEADER_PATTERN = re.compile(r"^(#{1,3})\s+(.+)", re.MULTILINE)

    text = Path(file_path).read_text(encoding="utf-8")
    file_name = Path(file_path).name

    matches = list(HEADER_PATTERN.finditer(text))
    if not matches:
        return [Document(
            page_content=text.strip(),
            metadata={"source": file_name, "file_type": "md"},
        )]

    docs = []
    for i, match in enumerate(matches):
        start = match.start()
        end = matches[i + 1].start() if i + 1 < len(matches) else len(text)
        section = text[start:end].strip()
        if len(section) < 20:
            continue
        docs.append(Document(
            page_content=section,
            metadata={
                "source": file_name,
                "section": match.group(2).strip(),
                "file_type": "md",
            },
        ))

    print(f"[MD] {file_name} → {len(docs)} 섹션 추출")
    return docs


# ─────────────────────────────────────────
# PPT 로더
# ─────────────────────────────────────────
def load_pptx(file_path: str) -> List[Document]:
    """
    python-pptx로 슬라이드별 텍스트 추출.
    텍스트 박스, 표 셀, 노트까지 수집합니다.
    """
    from pptx import Presentation

    docs = []
    prs = Presentation(file_path)

    for slide_num, slide in enumerate(prs.slides, start=1):
        texts = []

        # 텍스트 박스
        for shape in slide.shapes:
            if shape.has_text_frame:
                for para in shape.text_frame.paragraphs:
                    line = para.text.strip()
                    if line:
                        texts.append(line)

            # 표 셀
            if shape.has_table:
                for row in shape.table.rows:
                    for cell in row.cells:
                        cell_text = cell.text.strip()
                        if cell_text:
                            texts.append(cell_text)

        # 슬라이드 노트
        if slide.has_notes_slide:
            notes = slide.notes_slide.notes_text_frame.text.strip()
            if notes:
                texts.append(f"[노트] {notes}")

        if not texts:
            continue

        docs.append(Document(
            page_content="\n".join(texts),
            metadata={
                "source": Path(file_path).name,
                "slide": slide_num,
                "file_type": "pptx",
            }
        ))

    print(f"[PPT] {Path(file_path).name} → {len(docs)} 슬라이드 추출")
    return docs


# ─────────────────────────────────────────
# 통합 로더
# ─────────────────────────────────────────
def load_document(file_path: str) -> List[Document]:
    """확장자를 보고 적절한 로더를 자동 선택합니다."""
    ext = Path(file_path).suffix.lower()
    if ext == ".pdf":
        return load_pdf(file_path)
    elif ext in (".pptx", ".ppt"):
        return load_pptx(file_path)
    elif ext == ".md":
        return load_markdown(file_path)
    else:
        raise ValueError(f"지원하지 않는 파일 형식: {ext}")


SUPPORTED_EXTENSIONS = {".pdf", ".pptx", ".ppt", ".md"}


def load_directory(dir_path: str) -> List[Document]:
    """디렉터리 내 모든 PDF/PPT/MD 파일을 일괄 로드합니다."""
    all_docs = []
    for path in Path(dir_path).glob("**/*"):
        if path.suffix.lower() in SUPPORTED_EXTENSIONS:
            try:
                all_docs.extend(load_document(str(path)))
            except Exception as e:
                print(f"[WARN] {path.name} 로드 실패: {e}")
    print(f"\n총 {len(all_docs)}개 문서 로드 완료")
    return all_docs


# ─────────────────────────────────────────
# Adaptive 청킹 (메인)
# ─────────────────────────────────────────
MIN_CHUNK = 150   # 이보다 짧으면 다음 청크와 병합
MAX_CHUNK = 800   # 이보다 길면 재분할
OVERLAP   = 150


def chunk_adaptive(docs: List[Document]) -> List[Document]:
    """
    문서 구조에 맞는 adaptive 청킹.
    - MIN_CHUNK 미만: 인접 청크와 병합 (목차·짧은 헤더 제거)
    - MIN_CHUNK ~ MAX_CHUNK: 그대로 유지
    - MAX_CHUNK 초과: RecursiveCharacterTextSplitter로 재분할
    """
    from langchain_text_splitters import RecursiveCharacterTextSplitter

    splitter = RecursiveCharacterTextSplitter(
        chunk_size=MAX_CHUNK,
        chunk_overlap=OVERLAP,
        separators=["\n\n", "\n", "。", ". ", " ", ""],
    )

    # 1단계: 너무 작은 청크를 다음 청크에 병합
    merged: List[Document] = []
    buffer_text = ""
    buffer_meta: dict = {}

    for doc in docs:
        text = doc.page_content.strip()
        if not text:
            continue

        if not buffer_text:
            buffer_text = text
            buffer_meta = doc.metadata
        else:
            # 현재 버퍼가 MIN_CHUNK 미만이면 계속 합산
            if len(buffer_text) < MIN_CHUNK:
                buffer_text += "\n" + text
            else:
                merged.append(Document(page_content=buffer_text, metadata=buffer_meta))
                buffer_text = text
                buffer_meta = doc.metadata

    if buffer_text:
        merged.append(Document(page_content=buffer_text, metadata=buffer_meta))

    # 2단계: 너무 큰 청크를 재분할
    result: List[Document] = []
    for doc in merged:
        if len(doc.page_content) > MAX_CHUNK:
            splits = splitter.split_documents([doc])
            result.extend(splits)
        else:
            result.append(doc)

    sizes = [len(d.page_content) for d in result]
    avg = sum(sizes) // len(sizes) if sizes else 0
    print(f"adaptive 청킹: {len(docs)}개 → {len(result)}개 청크 (평균 {avg}자)")
    return result


# ─────────────────────────────────────────
# 레거시 청킹 (하위 호환)
# ─────────────────────────────────────────
def chunk_by_article(docs: List[Document]) -> List[Document]:
    """제X조 / 제X항 패턴 기반 청킹. 법령 형식 문서에 적합."""
    ARTICLE_PATTERN = re.compile(
        r"(?=제\s*\d+조|제\s*\d+항|\n\d+\.\s|\n[가-힣]\.\s)"
    )
    chunked = []
    for doc in docs:
        splits = ARTICLE_PATTERN.split(doc.page_content)
        for i, chunk in enumerate(splits):
            chunk = chunk.strip()
            if len(chunk) < 20:
                continue
            chunked.append(Document(page_content=chunk, metadata={**doc.metadata, "chunk_index": i}))
    print(f"조항 청킹: {len(docs)}개 → {len(chunked)}개 청크")
    return chunked


def chunk_by_size(
    docs: List[Document],
    chunk_size: int = 800,
    chunk_overlap: int = 150,
) -> List[Document]:
    """고정 크기 청킹. chunk_adaptive의 fallback."""
    from langchain_text_splitters import RecursiveCharacterTextSplitter

    splitter = RecursiveCharacterTextSplitter(
        chunk_size=chunk_size,
        chunk_overlap=chunk_overlap,
        separators=["\n\n", "\n", "。", ". ", " ", ""],
    )
    chunked = splitter.split_documents(docs)
    print(f"고정 크기 청킹: {len(docs)}개 → {len(chunked)}개 청크")
    return chunked
